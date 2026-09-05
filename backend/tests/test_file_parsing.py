"""Engineering-grade file parsers: every supported format returns a
ParsedDocument with the expected text/structure.

Fixtures are generated programmatically with the same libs the parsers use, so
these tests are self-contained (no on-disk sample files) and deterministic.
Heavy/optional behaviour (PDF scanned-page OCR) is exercised only for the
deterministic text-PDF path; OCR quality itself is a library concern covered
separately.
"""
from __future__ import annotations

import io
from pathlib import Path

import pytest

from app.rag.base import ParsedDocument
from app.rag.parsers import default_parser


# ---------------------------------------------------------------------------
# Fixture builders — write real (tiny) files in the supported binary formats.
# ---------------------------------------------------------------------------
def _write(path: Path, data: bytes) -> str:
    path.write_bytes(data)
    return str(path)


def _make_pdf(text: str) -> bytes:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), text)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def _make_docx(paras: list[str], table: list[list[str]] | None = None) -> bytes:
    import docx

    d = docx.Document()
    for p in paras:
        d.add_paragraph(p)
    if table:
        t = d.add_table(rows=len(table), cols=len(table[0]))
        for i, row in enumerate(table):
            for j, cell in enumerate(row):
                t.cell(i, j).text = cell
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _make_pptx(title: str, body: str) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(sheet: str, rows: list[list]) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = sheet
    for r in rows:
        ws.append(r)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_epub(title: str, body_html: str) -> bytes:
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_title(title)
    book.set_language("en")
    ch = epub.EpubHtml(title=title, file_name="ch.xhtml", content=body_html)
    book.add_item(ch)
    book.spine = [ch]
    # Required for ebooklib's reader to round-trip the spine/toc cleanly.
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    buf = io.BytesIO()
    epub.write_epub(buf, book, {})
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Plain-text family
# ---------------------------------------------------------------------------
def test_parse_txt(tmp_path: Path):
    p = _write(tmp_path / "n.txt", b"hello world")
    r = default_parser.parse(p, ".txt")
    assert isinstance(r, ParsedDocument)
    assert "hello world" in r.text
    assert r.metadata["parser_used"] == "text"


def test_parse_markdown(tmp_path: Path):
    p = _write(tmp_path / "n.md", b"# Title\n\nsome **bold** text")
    r = default_parser.parse(p, ".md")
    assert "Title" in r.text
    assert "bold" in r.text


def test_parse_json_as_text(tmp_path: Path):
    p = _write(tmp_path / "n.json", b'{"k": "v"}')
    r = default_parser.parse(p, ".json")
    assert '"k"' in r.text


# ---------------------------------------------------------------------------
# Office family
# ---------------------------------------------------------------------------
def test_parse_docx_with_table(tmp_path: Path):
    p = _write(tmp_path / "d.docx", _make_docx(["Alpha", "Beta"], [["a", "b"], ["c", "d"]]))
    r = default_parser.parse(p, ".docx")
    assert "Alpha" in r.text and "Beta" in r.text
    # Table folded into text + captured structurally.
    assert "a" in r.text and "d" in r.text
    assert r.tables and any("a" in t for t in r.tables)
    assert r.metadata["parser_used"] == "python-docx"


def test_parse_pptx(tmp_path: Path):
    p = _write(tmp_path / "s.pptx", _make_pptx("My Slide Title", "bullet content here"))
    r = default_parser.parse(p, ".pptx")
    assert "My Slide Title" in r.text
    assert "bullet content" in r.text
    assert r.metadata.get("slides") == 1
    assert r.pages and "My Slide Title" in r.pages[0]


def test_parse_xlsx(tmp_path: Path):
    p = _write(tmp_path / "t.xlsx", _make_xlsx("Data", [["name", "age"], ["alice", 30]]))
    r = default_parser.parse(p, ".xlsx")
    assert "alice" in r.text and "name" in r.text
    assert r.metadata.get("sheets") == ["Data"]


def test_parse_csv(tmp_path: Path):
    p = _write(tmp_path / "t.csv", b"name,age\nalice,30\n")
    r = default_parser.parse(p, ".csv")
    assert "alice" in r.text and "name" in r.text
    assert r.metadata["format"] == "csv"


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
def test_parse_pdf_text(tmp_path: Path):
    # Use >20 chars so the scanned-PDF OCR fallback (triggered on near-empty
    # text) does not kick in — we want the deterministic text-extraction path.
    p = _write(tmp_path / "d.pdf", _make_pdf("Hello PDF world, this is a longer test document."))
    r = default_parser.parse(p, ".pdf")
    assert "Hello PDF world" in r.text
    assert r.metadata.get("pages", 0) >= 1
    assert r.metadata["ocr_used"] is False


# ---------------------------------------------------------------------------
# Markup / e-book family
# ---------------------------------------------------------------------------
def test_parse_html(tmp_path: Path):
    html = b"<html><head><title>x</title></head><body><article><p>Hello HTML main</p></article></body></html>"
    p = _write(tmp_path / "p.html", html)
    r = default_parser.parse(p, ".html")
    assert "Hello HTML main" in r.text


def test_parse_rtf(tmp_path: Path):
    p = _write(tmp_path / "p.rtf", b"{\\rtf1\\ansi\\deff0 Hello RTF world}")
    r = default_parser.parse(p, ".rtf")
    assert "Hello RTF world" in r.text


def test_parse_epub(tmp_path: Path):
    p = _write(tmp_path / "b.epub", _make_epub("Book", "<html><body><p>Hello EPUB chapter</p></body></html>"))
    r = default_parser.parse(p, ".epub")
    assert "Hello EPUB chapter" in r.text
    assert r.metadata.get("chapters", 0) >= 1


# ---------------------------------------------------------------------------
# Dispatcher / contract
# ---------------------------------------------------------------------------
def test_parser_rejects_unsupported_and_images():
    with pytest.raises(ValueError):
        default_parser.parse("/tmp/x.unknown", ".unknown")
    # Images are intentionally not handled by the document parser.
    with pytest.raises(ValueError):
        default_parser.parse("/tmp/x.png", ".png")


def test_parseddocument_chars_property():
    pd = ParsedDocument(text="abc")
    assert pd.chars == 3
    assert ParsedDocument().chars == 0
