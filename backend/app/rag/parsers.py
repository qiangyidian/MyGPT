"""Document parsers.

One dispatcher (:class:`DefaultDocumentParser`) selects a parser by file
extension. Every parser returns a :class:`~app.rag.base.ParsedDocument`
(flat ``text`` + paginated ``pages`` + ``tables`` + parser ``metadata``).

Heavy/optional libs are imported lazily so the app still imports cleanly on a
minimal install; a missing parser surfaces a clear ``ValueError`` naming the
library that is missing. Images (png/jpg/...) are intentionally NOT mapped
here — they are handled by the multimodal/OCR path in ``attachment_service``,
which can pass them to a vision model directly or OCR them as a fallback.
"""
from __future__ import annotations

import logging
import os
import shutil
import subprocess
import tempfile
from typing import Any, Callable

from app.core.config import get_settings
from app.rag.base import DocumentParser, ParsedDocument

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Serialization helpers
# ---------------------------------------------------------------------------
def _table_to_text(rows: list[list[Any]]) -> str:
    """Render a 2-D cell matrix as ``col | col | …`` rows (embedding-friendly)."""
    if not rows:
        return ""
    out: list[str] = []
    for row in rows:
        cells = []
        for c in row:
            if c is None:
                cells.append("")
            else:
                cells.append(str(c).replace("\n", " ").replace("|", "/").strip())
        out.append(" | ".join(cells))
    return "\n".join(out)


def _df_to_text(df: Any, sheet: str | None = None) -> str:
    """Flatten a pandas DataFrame into the same ``col | col`` row format."""
    import pandas as pd  # lazy

    lines: list[str] = []
    if sheet:
        lines.append(f"## Sheet: {sheet}")
    cols = list(df.columns)
    lines.append(" | ".join(str(c) for c in cols))
    for _, row in df.iterrows():
        lines.append(" | ".join("" if pd.isna(v) else str(v) for v in row.tolist()))
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Plain-text family
# ---------------------------------------------------------------------------
def _parse_txt(path: str, _ext: str) -> ParsedDocument:
    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
        text = fh.read()
    return ParsedDocument(text=text, metadata={"parser_used": "text", "chars": len(text)})


def _parse_markdown(path: str, _ext: str) -> ParsedDocument:
    """Render Markdown to plain text (strip HTML tags from the rendered HTML)."""
    try:
        import markdown as md  # type: ignore
        from bs4 import BeautifulSoup  # type: ignore
    except Exception:
        # Without the libs, fall back to raw text — still usable for chunking.
        return _parse_txt(path, _ext)
    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
        html = md.markdown(fh.read())
    text = BeautifulSoup(html, "html.parser").get_text(separator="\n")
    return ParsedDocument(text=text, metadata={"parser_used": "markdown", "chars": len(text)})


# ---------------------------------------------------------------------------
# PDF (PyMuPDF primary: text + tables + scanned-page OCR fallback)
# ---------------------------------------------------------------------------
def _ocr_pdf_pages(doc: Any, max_pages: int) -> list[str]:
    """Render each page to an image and OCR it (scanned-PDF fallback)."""
    from app.rag import ocr

    pages: list[str] = []
    limit = min(len(doc), max(1, max_pages))
    for i in range(limit):
        try:
            page = doc[i]
            # Render at a modest DPI — 150 is a good accuracy/size tradeoff.
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("png")
            pages.append(ocr.image_to_text(img_bytes))
        except Exception as exc:  # noqa: BLE001 — one page must not kill the rest
            logger.debug("OCR page %d failed: %s", i, exc)
            pages.append("")
    return pages


def _parse_pdf(path: str, _ext: str) -> ParsedDocument:
    try:
        import pymupdf  # type: ignore  # PyMuPDF (importable as `fitz` too)
    except Exception as exc:  # pragma: no cover - optional dep
        raise ValueError(f"PyMuPDF 未安装，无法解析 PDF: {exc}") from exc

    settings = get_settings()
    pages: list[str] = []
    tables: list[str] = []
    ocr_used = False
    parser_used = "pymupdf"

    doc = pymupdf.open(path)
    try:
        for page in doc:
            text = (page.get_text("text") or "").strip()
            # Table extraction (PyMuPDF >= 1.23). Best-effort; some pages raise.
            try:
                finder = page.find_tables()
                for tbl in getattr(finder, "tables", []) or []:
                    rows = tbl.extract()
                    rendered = _table_to_text(rows)
                    if rendered.strip():
                        tables.append(rendered)
            except Exception:  # noqa: BLE001
                pass
            pages.append(text)
    finally:
        doc.close()

    full_text = "\n\n".join(p for p in pages if p)
    # Tables fold into the text so they are searchable / chunkable.
    if tables:
        full_text = (full_text + "\n\n" + "\n\n".join(tables)).strip()

    # Scanned-PDF fallback: near-empty text → OCR each rendered page.
    if settings.OCR_SCANNED_PDF and len(full_text.strip()) < 20:
        try:
            ocr_pages = _ocr_pdf_pages(pymupdf.open(path), settings.OCR_SCANNED_PDF_MAX_PAGES)
        except Exception as exc:  # noqa: BLE001
            logger.warning("scanned-PDF OCR failed: %s", exc)
            ocr_pages = []
        if any(p.strip() for p in ocr_pages):
            pages = ocr_pages
            full_text = "\n\n".join(p for p in ocr_pages if p.strip())
            ocr_used = True
            parser_used = "pymupdf+ocr"

    return ParsedDocument(
        text=full_text,
        pages=pages or None,
        tables=tables or None,
        metadata={
            "pages": len(pages),
            "tables": len(tables),
            "ocr_used": ocr_used,
            "parser_used": parser_used,
        },
    )


# ---------------------------------------------------------------------------
# Office family
# ---------------------------------------------------------------------------
def _parse_docx(path: str, _ext: str) -> ParsedDocument:
    try:
        import docx  # type: ignore  # python-docx
    except Exception as exc:  # pragma: no cover - optional dep
        raise ValueError(f"python-docx 未安装，无法解析 Word: {exc}") from exc
    d = docx.Document(path)
    parts: list[str] = [p.text for p in d.paragraphs if p.text and p.text.strip()]
    tables: list[str] = []
    for t in d.tables:
        rows = [[cell.text for cell in row.cells] for row in t.rows]
        rendered = _table_to_text(rows)
        if rendered.strip():
            tables.append(rendered)
    text = "\n".join(parts)
    if tables:
        text = (text + "\n\n" + "\n\n".join(tables)).strip()
    return ParsedDocument(
        text=text,
        tables=tables or None,
        metadata={"parser_used": "python-docx", "tables": len(tables)},
    )


def _parse_pptx(path: str, _ext: str) -> ParsedDocument:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dep
        raise ValueError(f"python-pptx 未安装，无法解析 PowerPoint: {exc}") from exc
    prs = Presentation(path)
    slides: list[str] = []
    tables: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        blocks: list[str] = [f"# 幻灯片 {idx}"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        blocks.append(txt)
            if getattr(shape, "has_table", False) and shape.has_table:
                tbl = shape.table
                rows = [[cell.text for cell in row.cells] for row in tbl.rows]
                rendered = _table_to_text(rows)
                if rendered.strip():
                    tables.append(rendered)
                    blocks.append(rendered)
        # Speaker notes (often carry the real explanation).
        if getattr(slide, "has_notes_slide", False) and slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                blocks.append(f"[备注] {notes}")
        slides.append("\n".join(blocks))
    text = "\n\n".join(slides)
    return ParsedDocument(
        text=text,
        pages=slides,
        tables=tables or None,
        metadata={"parser_used": "python-pptx", "slides": len(prs.slides), "tables": len(tables)},
    )


def _parse_table(path: str, ext: str) -> ParsedDocument:
    """Render CSV/XLSX/XLS into flattened text (per-sheet for workbooks)."""
    try:
        import pandas as pd  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dep
        raise ValueError(f"pandas 未安装，无法解析表格: {exc}") from exc

    if ext == ".csv":
        df = pd.read_csv(path)
        text = _df_to_text(df)
        return ParsedDocument(
            text=text,
            tables=[text],
            metadata={"format": "csv", "columns": list(df.columns), "rows": len(df),
                      "parser_used": "pandas"},
        )

    # Workbooks (.xlsx via openpyxl, legacy .xls via xlrd).
    engine = "xlrd" if ext == ".xls" else None
    xl = pd.ExcelFile(path, engine=engine) if engine else pd.ExcelFile(path)
    try:
        sheets = xl.sheet_names
        frames: list[str] = []
        for s in sheets:
            frames.append(_df_to_text(pd.read_excel(xl, sheet_name=s), sheet=s))
    finally:
        try:
            xl.close()
        except Exception:  # noqa: BLE001
            pass
    text = "\n\n".join(frames)
    return ParsedDocument(
        text=text,
        tables=frames,
        pages=frames,
        metadata={"format": ext.lstrip("."), "sheets": sheets, "parser_used": "pandas"},
    )


# ---------------------------------------------------------------------------
# Markup / e-book family
# ---------------------------------------------------------------------------
def _parse_html(path: str, _ext: str) -> ParsedDocument:
    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
        raw = fh.read()
    text = ""
    parser_used = "trafilatura"
    try:
        import trafilatura  # type: ignore

        text = trafilatura.extract(raw) or ""
    except Exception as exc:  # noqa: BLE001 — fall back to bs4
        logger.debug("trafilatura extract failed, using bs4: %s", exc)
        parser_used = "beautifulsoup"
    if not text.strip():
        from bs4 import BeautifulSoup  # type: ignore

        text = BeautifulSoup(raw, "html.parser").get_text(separator="\n")
        parser_used = "beautifulsoup"
    return ParsedDocument(text=text, metadata={"parser_used": parser_used, "chars": len(text)})


def _parse_epub(path: str, _ext: str) -> ParsedDocument:
    try:
        import ebooklib  # type: ignore
        from ebooklib import epub  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dep
        raise ValueError(f"EbookLib 未安装，无法解析 EPUB: {exc}") from exc
    from bs4 import BeautifulSoup  # type: ignore

    import warnings
    with warnings.catch_warnings():
        # ebooklib spams an EPUB-3 future-warning on every read; silence it.
        warnings.simplefilter("ignore")
        book = epub.read_epub(path)
    chapters: list[str] = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        txt = soup.get_text(separator="\n").strip()
        if txt:
            chapters.append(txt)
    text = "\n\n".join(chapters)
    return ParsedDocument(
        text=text,
        pages=chapters or None,
        metadata={"parser_used": "ebooklib", "chapters": len(chapters)},
    )


def _parse_rtf(path: str, _ext: str) -> ParsedDocument:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dep
        raise ValueError(f"striprtf 未安装，无法解析 RTF: {exc}") from exc
    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
        raw = fh.read()
    text = rtf_to_text(raw)
    return ParsedDocument(text=text, metadata={"parser_used": "striprtf", "chars": len(text)})


# ---------------------------------------------------------------------------
# Legacy .doc — python-docx cannot read the binary format; convert via
# LibreOffice when available, else surface a clear, actionable error.
# ---------------------------------------------------------------------------
def _find_soffice() -> str | None:
    return shutil.which("soffice") or shutil.which("libreoffice") or shutil.which("soffice.exe")


def _convert_via_libreoffice(path: str, target_ext: str, timeout: int = 120) -> str | None:
    soffice = _find_soffice()
    if not soffice:
        return None
    out_dir = tempfile.mkdtemp(prefix="lsoffice_")
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", target_ext.lstrip("."),
             "--outdir", out_dir, path],
            capture_output=True, timeout=timeout, check=False,
        )
        base = os.path.splitext(os.path.basename(path))[0]
        converted = os.path.join(out_dir, f"{base}{target_ext}")
        return converted if os.path.exists(converted) else None
    except Exception as exc:  # noqa: BLE001
        logger.debug("LibreOffice conversion failed: %s", exc)
        return None


def _parse_doc(path: str, _ext: str) -> ParsedDocument:
    converted = _convert_via_libreoffice(path, ".docx")
    if not converted:
        raise ValueError(
            "旧版 .doc（二进制 Word）需要 LibreOffice (soffice) 才能解析；"
            "请将文件另存为 .docx 后重新上传，或安装 LibreOffice 并确保 soffice 在 PATH。"
        )
    result = _parse_docx(converted, ".docx")
    result.metadata["parser_used"] = "libreoffice->python-docx"
    result.metadata["converted_from"] = ".doc"
    # Clean up the temp conversion (best-effort).
    try:
        os.remove(converted)
    except OSError:
        pass
    return result


def _parse_via_libreoffice(path: str, src_ext: str, target_ext: str, parse_fn):
    """Parse a mainstream-but-unsupported format by converting it via LibreOffice
    to an OOXML equivalent, then reusing the first-class parser.

    Covers OpenDocument (.odt/.ods/.odp) and legacy PowerPoint (.ppt). Requires
    LibreOffice (soffice) on PATH — same constraint as legacy .doc/.xls.
    """
    converted = _convert_via_libreoffice(path, target_ext)
    if not converted:
        raise ValueError(
            f"{src_ext} 需要 LibreOffice (soffice) 才能解析；"
            f"请将文件另存为 {target_ext} 后重新上传，或安装 LibreOffice 并确保 soffice 在 PATH。"
        )
    result = parse_fn(converted, target_ext)
    inner = result.metadata.get("parser_used", "?")
    result.metadata["parser_used"] = f"libreoffice->{inner}"
    result.metadata["converted_from"] = src_ext
    try:
        os.remove(converted)
    except OSError:
        pass
    return result


def _parse_odt(path: str, _ext: str) -> ParsedDocument:
    """OpenDocument Text (.odt) -> .docx via LibreOffice."""
    return _parse_via_libreoffice(path, ".odt", ".docx", _parse_docx)


def _parse_ods(path: str, _ext: str) -> ParsedDocument:
    """OpenDocument Spreadsheet (.ods) -> .xlsx via LibreOffice."""
    return _parse_via_libreoffice(path, ".ods", ".xlsx", _parse_table)


def _parse_odp(path: str, _ext: str) -> ParsedDocument:
    """OpenDocument Presentation (.odp) -> .pptx via LibreOffice."""
    return _parse_via_libreoffice(path, ".odp", ".pptx", _parse_pptx)


def _parse_ppt(path: str, _ext: str) -> ParsedDocument:
    """Legacy PowerPoint (.ppt) -> .pptx via LibreOffice."""
    return _parse_via_libreoffice(path, ".ppt", ".pptx", _parse_pptx)


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------
_EXT_TO_PARSER: dict[str, Callable[[str, str], ParsedDocument]] = {
    ".txt": _parse_txt,
    ".log": _parse_txt,
    ".json": _parse_txt,
    ".md": _parse_markdown,
    ".markdown": _parse_markdown,
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".doc": _parse_doc,
    ".pptx": _parse_pptx,
    ".csv": _parse_table,
    ".xlsx": _parse_table,
    ".xls": _parse_table,
    ".html": _parse_html,
    ".htm": _parse_html,
    ".epub": _parse_epub,
    ".rtf": _parse_rtf,
    # OpenDocument + legacy PowerPoint via LibreOffice conversion.
    ".odt": _parse_odt,
    ".ods": _parse_ods,
    ".odp": _parse_odp,
    ".ppt": _parse_ppt,
}


class DefaultDocumentParser(DocumentParser):
    """Dispatch parser by extension; raises ValueError on unsupported types.

    Images are intentionally unsupported here — they are owned by the
    multimodal/OCR path in ``attachment_service``.
    """

    def parse(self, file_path: str, file_type: str) -> ParsedDocument:
        ext = (file_type or "").lower()
        if not ext.startswith("."):
            ext = "." + ext
        fn = _EXT_TO_PARSER.get(ext)
        if fn is None:
            raise ValueError(f"不支持的文件类型: {ext}")
        return fn(file_path, ext)


default_parser = DefaultDocumentParser()
